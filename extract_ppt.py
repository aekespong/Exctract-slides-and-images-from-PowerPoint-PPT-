import os
import shutil
import time
from datetime import datetime

import win32com.client
from PIL import Image
from pptx import Presentation


def get_timestamped_dir(base_dir, name: str):
    """Skapar ett mappnamn med tidsstämpel"""
    timestamp = datetime.now().strftime("%y%m%d%H%M")
    return os.path.join(base_dir, f"{name}_{timestamp}")


def clean_directory(directory):
    """Tar bort och återskapar en mapp"""
    if os.path.exists(directory):
        shutil.rmtree(directory)
    os.makedirs(directory)


def extract_ppt_content(ppt_path, output_dir, name: str):
    """
    Extraherar text och bilder från en PowerPoint-fil.
    Sparar även varje slide som en bild.

    Args:
        ppt_path (str): Sökväg till PowerPoint-filen
        output_dir (str): Mapp där extraherat innehåll ska sparas
    """

    # clean_directory(output_dir)

    # Konvertera till absoluta sökvägar
    ppt_path = os.path.abspath(ppt_path)
    output_dir = get_timestamped_dir(os.path.abspath(output_dir), name)

    print(f"PowerPoint-fil: {ppt_path}")
    print(f"Output-mapp: {output_dir}")

    # Skapa output-mappar
    images_dir = os.path.join(output_dir, "images")
    slides_dir = output_dir  # os.path.join(output_dir, "slides")
    os.makedirs(images_dir, exist_ok=True)
    os.makedirs(slides_dir, exist_ok=True)

    # Öppna presentationen för textextrahering
    prs = Presentation(ppt_path)

    # Extrahera text och spara till fil
    text_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_content.append(f"=== Slide {i+1} ===\n" + "\n".join(slide_text))

    with open(
        os.path.join(output_dir, "extracted_text.txt"), "w", encoding="utf-8"
    ) as f:
        f.write("\n\n".join(text_content))

    print("Text har extraherats och sparats")

    # Extrahera inbäddade bilder
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 = bild
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext

                image_path = os.path.join(
                    images_dir, f"image_{image_count}.{image_format}"
                )
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                image_count += 1

    print(f"{image_count} bilder har extraherats")

    # Konvertera slides till bilder med PowerPoint COM
    powerpoint = None
    presentation = None

    try:
        print("Startar PowerPoint...")
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")

        print("Öppnar presentation...")
        presentation = powerpoint.Presentations.Open(ppt_path)

        # Skapa fullständig sökväg för output
        slides_output_path = os.path.join(slides_dir, "Slide")
        print(f"Förbereder att spara slides till: {slides_output_path}")

        # Spara varje slide som JPG
        presentation.SaveAs(
            slides_output_path,
            17,  # ppSaveAsJPG
        )

        print(f"Slides har sparats som bilder i: {slides_dir}")

    except Exception as e:
        print(f"Ett fel uppstod vid konvertering av slides: {str(e)}")
        raise

    finally:
        # Städa upp COM-objekt
        if presentation:
            print("Stänger presentation...")
            presentation.Close()
        if powerpoint:
            print("Avslutar PowerPoint...")
            powerpoint.Quit()
            time.sleep(1)  # Ge PowerPoint lite tid att stänga sig


# Exempel på användning
if __name__ == "__main__":
    # Ange fullständiga sökvägar till dina filer
    script_dir = os.path.dirname(os.path.abspath(__file__))
    ppt_path = os.path.join(script_dir, "Offert.pptx")
    output_dir = os.path.join(script_dir, "extracted")

    try:
        extract_ppt_content(ppt_path, output_dir, "Offert")
        print("Extrahering slutförd!")
    except Exception as e:
        print(f"Ett fel uppstod: {str(e)}")
